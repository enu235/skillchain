#!/usr/bin/env python3
"""Build the polished pptx deck for the TSLA / Claude Code research run."""

import json
import subprocess
import sys
from pathlib import Path

PYTHON_BIN = sys.executable

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --------------------------------------------------------------------------
# Paths
# --------------------------------------------------------------------------
RUN_DIR = Path("/Users/allan/dev/skillchain/research-output/run-20260503-210828-f48e")
DRILLDOWN_DIR = RUN_DIR / "drilldown"
CHARTS_DIR = RUN_DIR / "charts"
OUTLIERS_PATH = RUN_DIR / "outliers.json"
DRILLDOWN_INDEX_PATH = RUN_DIR / "drilldown-index.json"
OUTPUT_PATH = RUN_DIR / "deck.pptx"
RUN_ID = "run-20260503-210828-f48e"

CHART_RENDERER = Path(
    "/Users/allan/dev/skillchain/.claude/skills/chart-generation/scripts/render_chart.py"
)

CHARTS_DIR.mkdir(parents=True, exist_ok=True)

# --------------------------------------------------------------------------
# Style
# --------------------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x4E, 0x78)        # 31,78,120 title bar
NAVY_DARK = RGBColor(0x14, 0x36, 0x57)   # darker accent
INK = RGBColor(0x1A, 0x1A, 0x1A)         # body text
GRAY = RGBColor(0x66, 0x66, 0x66)        # caption gray
LIGHT_GRAY = RGBColor(0xC0, 0xC0, 0xC0)
ACCENT = RGBColor(0xC8, 0x5A, 0x2D)      # accent for callouts
SOFT_BG = RGBColor(0xF4, 0xF6, 0xF9)     # soft panel bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x2A, 0x3A)
CODE_FG = RGBColor(0xE6, 0xED, 0xF3)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TITLE_BAR_H = Inches(0.85)


# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------
def render_chart_png(spec, out_path):
    """Pipe a chart spec to render_chart.py and write a PNG."""
    proc = subprocess.run(
        [PYTHON_BIN, str(CHART_RENDERER), str(out_path)],
        input=json.dumps(spec),
        text=True,
        capture_output=True,
    )
    if proc.returncode != 0:
        raise RuntimeError(
            f"chart render failed for {out_path}: {proc.stderr.strip()}"
        )


def add_title_bar(slide, title_text):
    """Add the consistent navy title bar across the top."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, TITLE_BAR_H
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.shadow.inherit = False

    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"


def add_footer(slide, page_label=""):
    """Tiny footer with run id + optional page label."""
    foot = slide.shapes.add_textbox(
        Inches(0.4), SLIDE_H - Inches(0.35),
        SLIDE_W - Inches(0.8), Inches(0.28)
    )
    tf = foot.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = RUN_ID
    r.font.size = Pt(9)
    r.font.italic = True
    r.font.color.rgb = GRAY
    r.font.name = "Calibri"

    if page_label:
        # add right-aligned label on the same line via a second paragraph trick:
        # easier — second textbox on the right
        right = slide.shapes.add_textbox(
            SLIDE_W - Inches(2.5), SLIDE_H - Inches(0.35),
            Inches(2.1), Inches(0.28)
        )
        rtf = right.text_frame
        rtf.margin_left = 0
        rtf.margin_right = 0
        rtf.margin_top = 0
        rtf.margin_bottom = 0
        rp = rtf.paragraphs[0]
        rp.alignment = PP_ALIGN.RIGHT
        rr = rp.add_run()
        rr.text = page_label
        rr.font.size = Pt(9)
        rr.font.italic = True
        rr.font.color.rgb = GRAY
        rr.font.name = "Calibri"


def add_bullets(slide, left, top, width, height, bullets,
                font_size=16, line_spacing=1.15):
    """Add a list of bullet strings."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # leading bullet glyph
        bullet_run = p.add_run()
        bullet_run.text = "■  "  # filled square
        bullet_run.font.size = Pt(font_size)
        bullet_run.font.color.rgb = NAVY
        bullet_run.font.bold = True
        bullet_run.font.name = "Calibri"

        text_run = p.add_run()
        text_run.text = b
        text_run.font.size = Pt(font_size)
        text_run.font.color.rgb = INK
        text_run.font.name = "Calibri"
        p.space_after = Pt(6)


def add_caption(slide, left, top, width, text):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(11)
    r.font.italic = True
    r.font.color.rgb = GRAY
    r.font.name = "Calibri"


def add_section_label(slide, left, top, text, color=NAVY):
    tb = slide.shapes.add_textbox(left, top, Inches(6), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = "Calibri"


def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    """Add a dark code-style block with monospace font."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_size)
        r.font.name = "Consolas"
        r.font.color.rgb = CODE_FG


def add_panel(slide, left, top, width, height, fill=SOFT_BG, border=False):
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    panel.adjustments[0] = 0.04
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    if border:
        panel.line.color.rgb = LIGHT_GRAY
        panel.line.width = Pt(0.75)
    else:
        panel.line.fill.background()
    return panel


def add_stat_callout(slide, left, top, width, height, value, label,
                     value_color=NAVY):
    """Big-number + small label stat block."""
    add_panel(slide, left, top, width, height, fill=SOFT_BG, border=True)

    # value
    vbox = slide.shapes.add_textbox(
        left, top + Inches(0.15), width, Inches(0.9)
    )
    vtf = vbox.text_frame
    vtf.margin_left = 0
    vtf.margin_top = 0
    vp = vtf.paragraphs[0]
    vp.alignment = PP_ALIGN.CENTER
    vr = vp.add_run()
    vr.text = value
    vr.font.size = Pt(34)
    vr.font.bold = True
    vr.font.color.rgb = value_color
    vr.font.name = "Calibri"

    # label
    lbox = slide.shapes.add_textbox(
        left, top + Inches(1.05), width, Inches(0.7)
    )
    ltf = lbox.text_frame
    ltf.word_wrap = True
    ltf.margin_left = Inches(0.05)
    ltf.margin_top = 0
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    lr.text = label
    lr.font.size = Pt(11)
    lr.font.color.rgb = GRAY
    lr.font.name = "Calibri"


# --------------------------------------------------------------------------
# Load data
# --------------------------------------------------------------------------
with open(DRILLDOWN_INDEX_PATH) as f:
    drilldown_index = json.load(f)

drilldowns = {}
for entry in drilldown_index:
    with open(RUN_DIR / entry["path"]) as f:
        drilldowns[entry["slug"]] = json.load(f)

with open(OUTLIERS_PATH) as f:
    outliers = json.load(f)


# --------------------------------------------------------------------------
# Charts
# --------------------------------------------------------------------------
def build_chart_specs():
    """Choose the most informative chart per subtopic and build its spec."""
    specs = {}

    # claude-code-capabilities: Claude model pricing input vs output
    cc = drilldowns["claude-code-capabilities"]
    dp = cc["data_points"]
    specs["claude-code-capabilities"] = {
        "out": CHARTS_DIR / "claude-code-capabilities-pricing.png",
        "spec": {
            "title": "Claude model pricing per million tokens",
            "chart_type": "bar",
            "x_label": "Model",
            "y_label": "USD per million tokens",
            "series": [
                {"name": "Input",
                 "x": ["Opus 4.7", "Sonnet 4.6", "Haiku 4.5"],
                 "y": [dp[0]["value"], dp[2]["value"], dp[4]["value"]]},
                {"name": "Output",
                 "x": ["Opus 4.7", "Sonnet 4.6", "Haiku 4.5"],
                 "y": [dp[1]["value"], dp[3]["value"], dp[5]["value"]]},
            ],
        },
        "caption": "Source: platform.claude.com/docs/en/about-claude/pricing",
    }

    # claude-agent-sdk: Prompt cache multipliers
    sdk = drilldowns["claude-agent-sdk"]
    sdp = sdk["data_points"]
    specs["claude-agent-sdk"] = {
        "out": CHARTS_DIR / "claude-agent-sdk-cache.png",
        "spec": {
            "title": "Prompt-cache cost multipliers vs base input",
            "chart_type": "bar",
            "x_label": "Cache operation",
            "y_label": "Multiplier (x base input price)",
            "series": [{
                "x": ["5-min cache write",
                      "1-hour cache write",
                      "Cache read"],
                "y": [sdp[4]["value"], sdp[5]["value"], sdp[6]["value"]],
            }],
        },
        "caption": "Reads cost 0.1x base input; long-TTL writes cost 2x. Source: prompt-caching docs.",
    }

    # tsla-research-workflow: analyst price targets vs 52-week range
    tsla = drilldowns["tsla-research-workflow"]
    tdp = tsla["data_points"]
    specs["tsla-research-workflow"] = {
        "out": CHARTS_DIR / "tsla-research-workflow-targets.png",
        "spec": {
            "title": "TSLA analyst targets vs 52-week trading range (USD)",
            "chart_type": "bar",
            "x_label": "Reference point",
            "y_label": "USD per share",
            "series": [{
                "x": ["52-wk low", "52-wk high", "Median target",
                      "Bull (Wedbush)", "Bear (UBS)"],
                "y": [tdp[10]["value"], tdp[11]["value"], tdp[13]["value"],
                      tdp[14]["value"], tdp[15]["value"]],
            }],
        },
        "caption": "Bull/bear spread spans roughly +/-50% around the $398 median target.",
    }

    # local-llm-collaboration: VRAM for local models at Q4_K_M (8k ctx)
    llm = drilldowns["local-llm-collaboration"]
    ldp = llm["data_points"]
    specs["local-llm-collaboration"] = {
        "out": CHARTS_DIR / "local-llm-collaboration-vram.png",
        "spec": {
            "title": "VRAM required for local models @ Q4_K_M, 8k ctx",
            "chart_type": "bar",
            "x_label": "Model",
            "y_label": "VRAM (GB)",
            "series": [{
                "x": ["Llama 3.1 8B", "Llama 3.3 70B", "Qwen 2.5 72B",
                      "DeepSeek R1 14B", "Qwen 3 8B"],
                "y": [ldp[3]["value"], ldp[4]["value"], ldp[5]["value"],
                      ldp[6]["value"], ldp[7]["value"]],
            }],
        },
        "caption": "12.6x spread between 8B and 72B class — flagged as outlier.",
    }

    # hybrid-architecture: hallucination on ungrounded queries vs RAG-reduction
    hyb = drilldowns["hybrid-architecture"]
    # data_points[10] = 41% ungrounded; need a single chart with comparison.
    # We render a 2-bar comparison: ungrounded vs grounded estimate.
    specs["hybrid-architecture"] = {
        "out": CHARTS_DIR / "hybrid-architecture-hallucination.png",
        "spec": {
            "title": "Finance hallucination: ungrounded LLM vs RAG-grounded",
            "chart_type": "bar",
            "x_label": "Condition",
            "y_label": "Hallucination rate (%)",
            "series": [{
                "x": ["Ungrounded LLM", "Grounded (RAG, midpoint)",
                      "Grounded (RAG, best case)"],
                "y": [41.0, 41.0 * 0.30, 41.0 * 0.20],
            }],
        },
        "caption": "60-80% reduction with RAG/tool grounding (QuantMCP pattern).",
    }

    return specs


CHART_SPECS = build_chart_specs()
for slug, info in CHART_SPECS.items():
    render_chart_png(info["spec"], info["out"])
    print(f"rendered chart: {info['out']}")


# --------------------------------------------------------------------------
# Build the deck
# --------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]
slide_outline = []


def new_slide(title=None, page_label=None):
    slide = prs.slides.add_slide(BLANK)
    # white background is the pptx default; nothing to do.
    if title is not None:
        add_title_bar(slide, title)
    add_footer(slide, page_label or "")
    return slide


# --------------------------------------------------------------------------
# Slide 1 — Title
# --------------------------------------------------------------------------
def build_title_slide():
    slide = prs.slides.add_slide(BLANK)

    # full-bleed navy background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY

    # accent stripe on the left
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H
    )
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT

    # eyebrow
    eb = slide.shapes.add_textbox(
        Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.5)
    )
    ep = eb.text_frame.paragraphs[0]
    er = ep.add_run()
    er.text = "RESEARCH BRIEFING  ·  TSLA CASE STUDY"
    er.font.size = Pt(14)
    er.font.bold = True
    er.font.color.rgb = RGBColor(0xC8, 0xDC, 0xEB)
    er.font.name = "Calibri"

    # main title
    tb = slide.shapes.add_textbox(
        Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.6)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Deep Stock Research with"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Claude Code/SDK + a Local LLM"
    r2.font.size = Pt(40)
    r2.font.bold = True
    r2.font.color.rgb = WHITE
    r2.font.name = "Calibri"

    p3 = tf.add_paragraph()
    p3.space_before = Pt(10)
    r3 = p3.add_run()
    r3.text = "A TSLA case study for developers"
    r3.font.size = Pt(22)
    r3.font.italic = True
    r3.font.color.rgb = RGBColor(0xC8, 0xDC, 0xEB)
    r3.font.name = "Calibri"

    # divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(5.6),
        Inches(2.2), Emu(20000)
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

    # date + run id
    db = slide.shapes.add_textbox(
        Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.8)
    )
    dtf = db.text_frame
    dp = dtf.paragraphs[0]
    dr = dp.add_run()
    dr.text = "2026-05-03"
    dr.font.size = Pt(16)
    dr.font.color.rgb = WHITE
    dr.font.name = "Calibri"

    dp2 = dtf.add_paragraph()
    dr2 = dp2.add_run()
    dr2.text = f"Audience: developers building research workflows  ·  {RUN_ID}"
    dr2.font.size = Pt(11)
    dr2.font.italic = True
    dr2.font.color.rgb = RGBColor(0xA9, 0xC1, 0xD9)
    dr2.font.name = "Calibri"

    slide_outline.append("Title")


# --------------------------------------------------------------------------
# Slide 2 — Executive summary
# --------------------------------------------------------------------------
def build_exec_summary():
    slide = new_slide("Executive Summary", "2")
    bullets = [
        "Claude Code is a research-grade CLI: subagents, skills, hooks, MCP, plan mode, 1M-token context.",
        "The Agent SDK exposes the same loop programmatically (Python + TS); custom @tool MCP servers ship in-process.",
        "TSLA workflow has four stages: SEC EDGAR ingest, market data, news/sentiment, then frameworks and a bull/base/bear memo.",
        "Local LLMs (Ollama, LM Studio, llama.cpp) speak OpenAI- and Anthropic-compatible APIs — a 70B model fits in ~46 GB VRAM.",
        "Hybrid pattern: Claude Opus plans/critics, local 8B-22B does bulk work; MCP is the spine, files are the state.",
        "Anti-hallucination is non-negotiable in finance: tool-grounded RAG cuts the 41% ungrounded error rate by 60-80%.",
    ]
    add_bullets(
        slide, Inches(0.55), Inches(1.15),
        Inches(12.2), Inches(5.8), bullets, font_size=18
    )
    slide_outline.append("Executive Summary")


# --------------------------------------------------------------------------
# Slide 3 — Audience & scope
# --------------------------------------------------------------------------
def build_audience_scope():
    slide = new_slide("Audience & Scope", "3")

    # left column: who this is for
    add_section_label(slide, Inches(0.55), Inches(1.1),
                      "WHO THIS IS FOR")
    add_bullets(
        slide, Inches(0.55), Inches(1.5),
        Inches(6.2), Inches(5.5),
        [
            "Developers who already use Claude Code daily",
            "Quant-curious engineers exploring agentic research",
            "Builders running on-device LLMs alongside cloud APIs",
            "Teams cost-optimizing premium reasoning vs bulk work",
        ],
        font_size=16,
    )

    # right column: scope
    add_section_label(slide, Inches(7.0), Inches(1.1),
                      "WHAT WE COVER")
    add_bullets(
        slide, Inches(7.0), Inches(1.5),
        Inches(5.8), Inches(5.5),
        [
            "Claude Code CLI surfaces: subagents, skills, hooks, MCP",
            "Claude Agent SDK: query(), ClaudeSDKClient, custom tools",
            "End-to-end TSLA research pipeline (sources -> memo)",
            "Pairing with Ollama / LM Studio / llama.cpp",
            "Hybrid topology, pitfalls, and concrete next steps",
        ],
        font_size=16,
    )
    slide_outline.append("Audience & Scope")


# --------------------------------------------------------------------------
# Slide 4 — Architecture overview
# --------------------------------------------------------------------------
def build_architecture():
    slide = new_slide("Hybrid Architecture at a Glance", "4")

    # 4 role cards across the top
    roles = [
        ("PLANNER / CRITIC",
         "Claude Opus 4.7 (1M ctx)",
         "Decomposes the topic, writes plans, reviews drafts, makes the final call."),
        ("BULK WORKER",
         "Local LLM (8B-22B, Ollama)",
         "Summarizes filings, drafts sections, runs overnight iterations cheaply."),
        ("INTEGRATION SPINE",
         "Model Context Protocol",
         "MCP servers expose tools (SEC, prices, browser, local LLM) to both tiers."),
        ("STATE STORE",
         "File artifacts + SQLite",
         "JSON drilldowns, xlsx workbooks, charts — durable handoffs between runs."),
    ]
    card_w = Inches(3.05)
    card_h = Inches(2.2)
    gap = Inches(0.12)
    left0 = Inches(0.5)
    top = Inches(1.2)
    for i, (header, sub, body) in enumerate(roles):
        x = left0 + (card_w + gap) * i
        add_panel(slide, x, top, card_w, card_h, fill=SOFT_BG, border=True)

        # accent strip on left
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, top, Inches(0.08), card_h
        )
        strip.line.fill.background()
        strip.fill.solid()
        strip.fill.fore_color.rgb = NAVY

        tb = slide.shapes.add_textbox(
            x + Inches(0.2), top + Inches(0.15),
            card_w - Inches(0.3), card_h - Inches(0.25)
        )
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = header
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = ACCENT
        r.font.name = "Calibri"

        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(15)
        r2.font.bold = True
        r2.font.color.rgb = NAVY
        r2.font.name = "Calibri"

        p3 = tf.add_paragraph()
        p3.space_before = Pt(6)
        r3 = p3.add_run()
        r3.text = body
        r3.font.size = Pt(11)
        r3.font.color.rgb = INK
        r3.font.name = "Calibri"

    # bottom: data-flow line as text
    add_section_label(slide, Inches(0.55), Inches(3.7), "DATA FLOW")
    add_bullets(
        slide, Inches(0.55), Inches(4.1),
        Inches(12.3), Inches(2.6),
        [
            "Prompt -> Claude planner decomposes into subtopics; writes drilldown specs.",
            "Subagents fan out: each calls MCP tools (WebFetch, SEC, Polygon, local-LLM).",
            "Local LLM does the bulk read-and-summarize; returns JSON to planner.",
            "Planner critiques, runs outlier detection, then assembles xlsx + pptx artifacts.",
            "Hooks (PreToolUse/PostToolUse) gate tool calls; /loop schedules recurring runs.",
        ],
        font_size=15,
    )
    slide_outline.append("Hybrid Architecture at a Glance")


# --------------------------------------------------------------------------
# Generic two-slide subtopic builder
# --------------------------------------------------------------------------
def build_subtopic(slug, narrative_title, narrative_bullets,
                   data_title, chart_caption=None,
                   code_lines=None, code_label=None,
                   stat_callouts=None):
    """
    narrative slide on the left + data slide with embedded chart.
    stat_callouts: optional list of (value, label) tuples for the data slide.
    code_lines: optional list[str] for an embedded code block on the narrative slide.
    """
    # ---- narrative slide ----
    slide = new_slide(narrative_title)

    if code_lines:
        # bullets on left, code on right
        add_bullets(
            slide, Inches(0.55), Inches(1.2),
            Inches(7.4), Inches(5.6), narrative_bullets, font_size=15
        )
        if code_label:
            add_section_label(slide, Inches(8.2), Inches(1.2), code_label)
            code_top = Inches(1.55)
        else:
            code_top = Inches(1.2)
        add_code_block(
            slide, Inches(8.2), code_top,
            Inches(4.7), Inches(4.6),
            code_lines, font_size=11
        )
    else:
        add_bullets(
            slide, Inches(0.55), Inches(1.2),
            Inches(12.3), Inches(5.6), narrative_bullets, font_size=17
        )

    slide_outline.append(narrative_title)

    # ---- data slide ----
    slide2 = new_slide(data_title)

    if stat_callouts:
        # row of stat callouts at top, chart below
        n = len(stat_callouts)
        block_w = Inches(2.55)
        gap = Inches(0.18)
        total_w = block_w * n + gap * (n - 1)
        start_x = (SLIDE_W - total_w) // 2
        top_y = Inches(1.15)
        for i, (val, label) in enumerate(stat_callouts):
            x = start_x + (block_w + gap) * i
            add_stat_callout(
                slide2, x, top_y, block_w, Inches(1.55),
                val, label
            )
        chart_top = Inches(2.85)
        chart_h = Inches(4.0)
    else:
        chart_top = Inches(1.15)
        chart_h = Inches(5.4)

    chart_path = CHART_SPECS[slug]["out"]
    # 1280x720 -> aspect 16:9, fit width 11", center.
    img_w = Inches(11.0)
    img_h = Inches(6.18)
    if img_h > chart_h:
        img_h = chart_h
        img_w = img_h * (16 / 9)
    img_left = (SLIDE_W - img_w) // 2
    slide2.shapes.add_picture(
        str(chart_path), img_left, chart_top, width=img_w, height=img_h
    )

    cap = chart_caption or CHART_SPECS[slug]["caption"]
    add_caption(
        slide2, Inches(0.6),
        chart_top + img_h + Inches(0.05),
        Inches(12.2), cap
    )
    slide_outline.append(data_title)


# --------------------------------------------------------------------------
# Slides 5-6 — Claude Code capabilities
# --------------------------------------------------------------------------
def build_claude_code():
    bullets = [
        "Subagents isolate context: own tools, prompt, summary-only return.",
        "Built-ins: Explore, Plan, general-purpose. Custom in .claude/agents/.",
        "Skills at .claude/skills/<name>/SKILL.md with YAML frontmatter.",
        "Hooks fire on lifecycle events from settings.json (PreToolUse, etc.).",
        "MCP via 'claude mcp add' over HTTP / stdio / SSE — three scopes.",
        "Plan mode = read-only tools + reviewable plan before any side effect.",
    ]
    code = [
        "# wire an MCP server into the project scope",
        "claude mcp add --scope project sec-edgar \\",
        "    --transport stdio \\",
        "    -- python servers/sec_edgar.py",
        "",
        "# verify it landed and is reachable",
        "claude mcp list",
        "",
        "# inside a session, scope a one-off",
        "/permissions allow Bash:claude mcp call sec-edgar:get_filing",
    ]
    build_subtopic(
        slug="claude-code-capabilities",
        narrative_title="Claude Code: research-grade CLI primitives",
        narrative_bullets=bullets,
        data_title="Claude Code: model pricing per 1M tokens",
        code_lines=code,
        code_label="MCP wiring (claude mcp add)",
        stat_callouts=[
            ("1M", "context window @ standard pricing"),
            ("0.1x", "cache-read multiplier"),
            ("$10", "per 1,000 web searches"),
        ],
    )


# --------------------------------------------------------------------------
# Slides 7-8 — Claude Agent SDK
# --------------------------------------------------------------------------
def build_agent_sdk():
    bullets = [
        "Official Python (claude-agent-sdk) and TS (@anthropic-ai/claude-agent-sdk).",
        "query() = single-shot async iterator; ClaudeSDKClient = bidirectional.",
        "Custom tools via @tool + create_sdk_mcp_server (in-process MCP).",
        "Streaming opt-in with include_partial_messages -> StreamEvent stream.",
        "Sub-agents via AgentDefinition delegate focused work with own toolsets.",
        "Headless mode + Bedrock/Vertex/Foundry providers ready for CI/CD.",
    ]
    code = [
        "import anyio",
        "from claude_agent_sdk import query, ClaudeAgentOptions",
        "",
        "async def main():",
        "    opts = ClaudeAgentOptions(",
        "        model='claude-opus-4-7',",
        "        system_prompt='You are a TSLA research analyst.',",
        "        allowed_tools=['WebFetch', 'mcp__sec__get_filing'],",
        "    )",
        "    async for msg in query(prompt='Brief TSLA Q1 2026', options=opts):",
        "        print(msg)",
        "",
        "anyio.run(main)",
    ]
    build_subtopic(
        slug="claude-agent-sdk",
        narrative_title="Claude Agent SDK: programmatic research loop",
        narrative_bullets=bullets,
        data_title="Agent SDK: prompt-cache cost multipliers",
        code_lines=code,
        code_label="query() with custom tools (Python)",
        stat_callouts=[
            ("0.1.72", "Python SDK current"),
            ("0.2.111+", "min for Opus 4.7"),
            ("50%", "Batch API discount"),
        ],
    )


# --------------------------------------------------------------------------
# Slides 9-10 — TSLA workflow
# --------------------------------------------------------------------------
def build_tsla_workflow():
    bullets = [
        "Stage 1 - Primary disclosures: SEC EDGAR (10-K/10-Q/8-K) + Tesla IR.",
        "Stage 2 - Market data: Polygon (prod), Tiingo (EOD), yfinance (proto).",
        "Stage 3 - News & sentiment: Finnhub, NewsAPI, social-media scrapers.",
        "Stage 4 - Frameworks: fundamentals, technicals, macro, narrative.",
        "Outputs: exec summary, price-vs-catalyst chart, fundamentals table, bull/base/bear memo.",
        "Event-trigger re-runs on new 8-K filings (e.g., Q1 2026 4/22, deliveries 4/2).",
    ]
    build_subtopic(
        slug="tsla-research-workflow",
        narrative_title="TSLA workflow: from EDGAR to memo",
        narrative_bullets=bullets,
        data_title="TSLA: analyst-target spread and trading range",
        stat_callouts=[
            ("$22.38B", "Q1 2026 revenue (+16% YoY)"),
            ("21.1%", "Q1 2026 gross margin"),
            ("357", "P/E ratio (May 1 2026)"),
        ],
    )


# --------------------------------------------------------------------------
# Slides 11-12 — Local LLM collaboration
# --------------------------------------------------------------------------
def build_local_llm():
    bullets = [
        "Ollama @ :11434, LM Studio @ :1234, llama-server @ :8080.",
        "All speak OpenAI-compat; llama.cpp also speaks Anthropic /v1/messages.",
        "Point Claude Code at Ollama: ANTHROPIC_BASE_URL=http://localhost:11434.",
        "OllamaClaude MCP cuts Anthropic tokens up to ~98.75% on file work.",
        "Routing: ~70% bulk to local/open, ~10% premium reasoning to Claude.",
        "Q4_K_M sizing: 8B ~6 GB, 14B ~11 GB, 70B ~46 GB, 72B ~50 GB.",
    ]
    code = [
        "# pull a fast worker model",
        "ollama pull llama3.1:8b",
        "",
        "# enable Anthropic-compatible mode",
        "export ANTHROPIC_BASE_URL=http://localhost:11434",
        "export ANTHROPIC_AUTH_TOKEN=ollama",
        "export ANTHROPIC_API_KEY=''",
        "",
        "# Claude Code now uses Ollama as the backend",
        "claude --model llama3.1:8b",
    ]
    build_subtopic(
        slug="local-llm-collaboration",
        narrative_title="Local LLMs: Ollama, LM Studio, llama.cpp",
        narrative_bullets=bullets,
        data_title="Local LLM VRAM footprint (Q4_K_M, 8k ctx)",
        code_lines=code,
        code_label="Wire Claude Code to Ollama",
        stat_callouts=[
            ("~6 GB", "Llama 3.1 8B VRAM"),
            ("~46 GB", "Llama 3.3 70B VRAM"),
            ("98.75%", "tokens saved (file ops)"),
        ],
    )


# --------------------------------------------------------------------------
# Slides 13-14 — Hybrid architecture
# --------------------------------------------------------------------------
def build_hybrid():
    bullets = [
        "MCP spec 2025-11-25: async Tasks, server-side loops, parallel calls.",
        "LiteLLM: one OpenAI-compatible gateway in front of 100+ providers.",
        "Opus 4.7 (1M ctx) plans/critiques; Llama 3.1 8B (128k) does bulk work.",
        "Watch out: Ollama default ctx is 2048 - silent truncation in pipelines.",
        "Hooks gate tool calls; headless 'claude -p' fits GitHub Actions / cron.",
        "Anti-hallucination: never read finance facts from parametric memory.",
    ]
    code = [
        "# .claude/settings.json (excerpt)",
        '{',
        '  "hooks": {',
        '    "PostToolUse": [{',
        '      "matcher": {"tool": "Bash"},',
        '      "command": "scripts/audit_tool_call.py"',
        '    }]',
        '  }',
        '}',
        "",
        "# headless run, exits on completion",
        "claude -p 'Run TSLA daily brief; write to artifacts/'",
    ]
    build_subtopic(
        slug="hybrid-architecture",
        narrative_title="Hybrid topology: planner, worker, spine, state",
        narrative_bullets=bullets,
        data_title="Why grounding matters: finance hallucination rate",
        code_lines=code,
        code_label="Hooks + headless run",
        stat_callouts=[
            ("41%", "ungrounded finance error rate"),
            ("60-80%", "reduction with RAG/tools"),
            ("64.4%", "Opus 4.7 Finance Agent v1.1"),
        ],
    )


# --------------------------------------------------------------------------
# Slide 15 — Outliers / risks
# --------------------------------------------------------------------------
def build_outliers():
    slide = new_slide("Outliers, contradictions & risks", "15")

    sev_color = {
        "high":   RGBColor(0xB0, 0x2A, 0x2A),
        "medium": ACCENT,
        "low":    RGBColor(0x4F, 0x6F, 0x8F),
    }

    anomalies = outliers["anomalies"]
    # Simple stacked rows
    top = Inches(1.15)
    row_h = Inches(1.45)
    row_gap = Inches(0.05)

    for i, an in enumerate(anomalies):
        y = top + (row_h + row_gap) * i
        # row panel
        add_panel(
            slide, Inches(0.5), y,
            SLIDE_W - Inches(1.0), row_h,
            fill=SOFT_BG, border=True
        )
        # severity pill on left
        sev = an["severity"]
        pill_color = sev_color.get(sev, GRAY)
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), y + Inches(0.22),
            Inches(1.4), Inches(0.45)
        )
        pill.adjustments[0] = 0.5
        pill.line.fill.background()
        pill.fill.solid()
        pill.fill.fore_color.rgb = pill_color
        ptf = pill.text_frame
        ptf.margin_left = Inches(0.05)
        ptf.margin_right = Inches(0.05)
        ptf.margin_top = Inches(0.03)
        ptf.margin_bottom = Inches(0.03)
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = sev.upper()
        pr.font.size = Pt(10)
        pr.font.bold = True
        pr.font.color.rgb = WHITE
        pr.font.name = "Calibri"

        # type label below pill
        tlb = slide.shapes.add_textbox(
            Inches(0.7), y + Inches(0.7),
            Inches(1.5), Inches(0.4)
        )
        ttp = tlb.text_frame.paragraphs[0]
        ttp.alignment = PP_ALIGN.CENTER
        ttr = ttp.add_run()
        ttr.text = an["type"].replace("_", " ")
        ttr.font.size = Pt(9)
        ttr.font.italic = True
        ttr.font.color.rgb = GRAY
        ttr.font.name = "Calibri"

        # body
        body = slide.shapes.add_textbox(
            Inches(2.3), y + Inches(0.1),
            Inches(10.5), row_h - Inches(0.2)
        )
        btf = body.text_frame
        btf.word_wrap = True
        btf.margin_left = 0
        btf.margin_top = 0

        # subtopic refs as a small header
        bp = btf.paragraphs[0]
        br = bp.add_run()
        br.text = "  ·  ".join(an["involved_subtopics"])
        br.font.size = Pt(11)
        br.font.bold = True
        br.font.color.rgb = NAVY
        br.font.name = "Calibri"

        bp2 = btf.add_paragraph()
        bp2.space_before = Pt(2)
        # truncate description to a single readable bullet
        desc = an["description"]
        # keep first ~340 chars to fit row
        if len(desc) > 360:
            desc = desc[:357].rstrip() + "..."
        br2 = bp2.add_run()
        br2.text = desc
        br2.font.size = Pt(11)
        br2.font.color.rgb = INK
        br2.font.name = "Calibri"

    slide_outline.append("Outliers, contradictions & risks")


# --------------------------------------------------------------------------
# Slide 16 — How to start / next steps
# --------------------------------------------------------------------------
def build_how_to_start():
    slide = new_slide("How to start (concrete next steps)", "16")

    # left: numbered steps
    add_section_label(slide, Inches(0.55), Inches(1.1), "FIRST 30 MINUTES")
    steps = [
        "Install Claude Code; verify with 'claude --version'.",
        "Pull a local worker model with Ollama.",
        "Wire SEC EDGAR + a market-data MCP server.",
        "Add a logging hook in settings.json.",
        "Run a TSLA brief in headless mode and inspect artifacts.",
        "Schedule a /loop or cron run for daily refresh.",
    ]
    add_bullets(
        slide, Inches(0.55), Inches(1.5),
        Inches(6.4), Inches(5.4),
        steps, font_size=15
    )

    # right: copy-paste commands as a code block
    add_section_label(slide, Inches(7.2), Inches(1.1), "COMMANDS")
    cmds = [
        "# 1. Install the SDK",
        "pip install claude-agent-sdk",
        "",
        "# 2. Pull a worker model",
        "ollama pull llama3.1:8b",
        "",
        "# 3. Wire an MCP server",
        "claude mcp add --scope project sec-edgar \\",
        "    --transport stdio -- python sec.py",
        "",
        "# 4. Headless brief",
        "claude -p 'Brief TSLA Q1 2026; bull/base/bear' \\",
        "    --output-format json > brief.json",
        "",
        "# 5. Recurring run",
        "claude /loop 24h /tsla-brief",
    ]
    add_code_block(
        slide, Inches(7.2), Inches(1.5),
        Inches(5.6), Inches(5.4),
        cmds, font_size=11
    )
    slide_outline.append("How to start (concrete next steps)")


# --------------------------------------------------------------------------
# Slide 17 — Sources & references
# --------------------------------------------------------------------------
def build_sources():
    slide = new_slide("Sources & references", "17")

    # Collect deduped sources
    seen = set()
    grouped = []
    for slug in ["claude-code-capabilities", "claude-agent-sdk",
                 "tsla-research-workflow", "local-llm-collaboration",
                 "hybrid-architecture"]:
        urls = []
        for u in drilldowns[slug]["sources"]:
            if u not in seen:
                seen.add(u)
                urls.append(u)
        grouped.append((slug, urls))

    # Two columns of URLs
    col_w = Inches(6.0)
    left_x = [Inches(0.55), Inches(7.0)]
    top = Inches(1.15)
    col_h = Inches(5.7)

    # interleave subtopics into 2 columns by cumulative line count
    col_lines = [[], []]
    col_count = [0, 0]
    for slug, urls in grouped:
        # decide which column (shorter so far)
        c = 0 if col_count[0] <= col_count[1] else 1
        col_lines[c].append(("__header__", slug))
        col_count[c] += 1
        for u in urls:
            col_lines[c].append(("url", u))
            col_count[c] += 1
        # spacer
        col_lines[c].append(("__spacer__", ""))
        col_count[c] += 1

    for ci in range(2):
        tb = slide.shapes.add_textbox(left_x[ci], top, col_w, col_h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_top = 0

        first = True
        for kind, val in col_lines[ci]:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            if kind == "__header__":
                r = p.add_run()
                r.text = val
                r.font.size = Pt(12)
                r.font.bold = True
                r.font.color.rgb = NAVY
                r.font.name = "Calibri"
                p.space_before = Pt(2)
            elif kind == "__spacer__":
                r = p.add_run()
                r.text = ""
                r.font.size = Pt(4)
            else:
                r = p.add_run()
                r.text = val
                r.font.size = Pt(9)
                r.font.color.rgb = INK
                r.font.name = "Consolas"
                p.space_after = Pt(1)

    add_caption(
        slide, Inches(0.55), Inches(6.85), Inches(12.2),
        f"Full machine-readable corpus: drilldown/*.json, outliers.json, report.xlsx in {RUN_ID}"
    )
    slide_outline.append("Sources & references")


# --------------------------------------------------------------------------
# Build all slides in order
# --------------------------------------------------------------------------
build_title_slide()        # 1
build_exec_summary()       # 2
build_audience_scope()     # 3
build_architecture()       # 4
build_claude_code()        # 5, 6
build_agent_sdk()          # 7, 8
build_tsla_workflow()      # 9, 10
build_local_llm()          # 11, 12
build_hybrid()             # 13, 14
build_outliers()           # 15
build_how_to_start()       # 16
build_sources()            # 17

prs.save(str(OUTPUT_PATH))
print(f"\nDeck saved to {OUTPUT_PATH}")
print(f"Slide count: {len(slide_outline)}")
for i, t in enumerate(slide_outline, 1):
    print(f"  {i:2d}. {t}")
